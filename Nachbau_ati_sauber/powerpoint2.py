from pptx import Presentation
from pptx.util import Inches
import os

def erstelle_powerpoint(bilder_liste, speichername="Präsentation.pptx"):
    # Erstelle eine neue Präsentation
    prs = Presentation()

    # Foliengröße anpassen (Standard ist 16:9, hier optional 4:3)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for bild in bilder_liste:
        # Erstelle eine Folie mit leerem Layout
        slide_layout = prs.slide_layouts[6]  # Layout ohne Titel oder Text
        slide = prs.slides.add_slide(slide_layout)

        # Lade das Bild
        img_path = bild

        # Berechne die Bildgröße für die gesamte Folie
        left = 0
        top = 0
        width = prs.slide_width
        height = prs.slide_height

        # Füge das Bild zur Folie hinzu
        slide.shapes.add_picture(img_path, left, top, width, height)

    # Speichere die Präsentation
    prs.save(speichername)
    print(f"PowerPoint gespeichert als: {speichername}")

# Liste mit Bildern (ersetze diese Pfade durch deine eigenen)
bilder_liste = [
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\Original_ohne_LTF.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\Original.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\0.95Volt.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\0.6char.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt_0.6char.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt_0.6char_4al.PNG",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt_0.6char_4al_kontakt30.PNG"
]


bilder_liste = [
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\charzucont1.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\charzucont2.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Einfallswinkelalpha_19-21.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Einfallswinkelalpha_15-25.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Emax_39-41.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Emax.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\sigma.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\activeLayer.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\activeLayer_2-4.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Totschicht_0.02-0.08.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Totschicht_2-8_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Kontaktmaterialdicke_49-51.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Kontaktmaterialdicke_49-51_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Kontaktmaterialdicke.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Kontaktmaterialdicke_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Bedeckungsfaktor.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation\\Bedeckungsfaktor_zoom.png"

]



bilder_liste = [
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\activeLayer_2.9-3.1.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\activeLayer_2-4.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\charzucont_1.2-0.8.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\charzucont_1.2-0.8_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Einfallswinkelalpha_15-25.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Einfallswinkelalpha_19-21.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Emax_39-41.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Emax_39-41_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Kontaktmaterialdicke_40-60.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Kontaktmaterialdicke_40-60_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\sigma_0.8-1.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\sigma_0.8-1_zoom.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\sigma_1.0214-1.0414.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Totschicht_0.02-0.08.png",
    "C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation_2\\Parametervariation_2\\Totschicht_0.02-0.08_zoom.png"




]


# PowerPoint erstellen
erstelle_powerpoint(bilder_liste, "Meine_Praesentation.pptx")

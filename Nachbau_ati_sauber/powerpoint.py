from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

# Erstelle ein PowerPoint-Objekt
prs = Presentation()

def Folie(Bild1, Bild2, Text):

    # Füge eine Folie hinzu
    slide_layout = prs.slide_layouts[6]  # Layout ohne Titel
    slide = prs.slides.add_slide(slide_layout)

    # Füge das erste Bild ein
    left = Inches(0)  # Position von links
    top = Inches(0)     # Position von oben
    width = Inches(7)   # Breite des Bildes
    height = Inches(4.3)
    img_path1 = Bild1   # Pfad zum ersten Bild
    slide.shapes.add_picture(img_path1, left, top, width, height)

    # Füge das zweite Bild ein
    top = Inches(3) # Position von oben, um das Bild untereinander zu platzieren
    img_path2 = Bild2   # Pfad zum zweiten Bild
    slide.shapes.add_picture(img_path2, left, top , width, height)

    # Beispiel: Einfügen eines Textfeldes
    left = Inches(7.5)  # Position von links
    top = Inches(2.6)  # Position von oben
    width = Inches(2)  # Breite des Textfeldes
    height = Inches(4)  # Höhe des Textfeldes


    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame

    # Füge einen Absatz hinzu und formatiere ihn
    paragraph = text_frame.add_paragraph()
    paragraph.text = Text
    paragraph.font.bold = True  # Der gesamte Absatz wird fett
    paragraph.font.size = Pt(24)  # Schriftgröße
    #paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Textfarbe rot

    # Erster Pfeil (horizontal)
    start_x = Inches(8.5)
    start_y = Inches(2.8)
    end_x = Inches(7)  # Ende des Pfeils
    end_y = Inches(1.6)

    arrow = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )
    line_format = arrow.line
    line_format.color.rgb = RGBColor(0, 0, 0)  # Schwarz
    line_format.width = Inches(0.1)
    line_format.end_arrowhead = True  # Pfeilspitze hinzufügen

    # Zweiter Pfeil (vertikal)
    start_x = Inches(8.5)
    start_y = Inches(3.9)
    end_x = Inches(7)
    end_y = Inches(5.1)

    arrow2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )
    line_format = arrow2.line
    line_format.color.rgb = RGBColor(0, 255, 0)  # Grün
    line_format.width = Inches(0.1)
    line_format.end_arrowhead = True  # Pfeilspitze hinzufügen


B1 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\Original_ohne_LTF.PNG'
B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\Original.PNG'
T = "mit Lochtransfer-\nfaktor"
Folie(B1,B2,T)


B1 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\Original.PNG'
B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\0.95Volt.PNG'
T = "Spannung\n    -5%"
Folie(B1,B2,T)

B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt.PNG'
T = "Spannung\n   -2kV"
Folie(B1,B2,T)

B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\0.6char.PNG'
T = "Charakteristische\nStrahlung*0,6"
Folie(B1,B2,T)

B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt_0.6char.PNG'
T = "char*0,6\n-2kV"
Folie(B1,B2,T)

B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt_0.6char_4al.PNG'
T = "char*0.6,-2kV\nactivelayer 4[mm]"
Folie(B1,B2,T)

B2 = 'C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\neue Plots\\-2Volt_0.6char_4al_kontakt30.PNG'
T = "char*0.6,-2kV,al 4mm\ncontaktlayer 30nm"
Folie(B1,B2,T)




# Speichere die PowerPoint-Datei
prs.save('C:\\Users\\julia\\OneDrive\\Dokumente\\A_Christian\\Masterarbeit\\Präsentation\\Bilder.pptx')
